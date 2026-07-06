from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

from .extract_pdf import extract_pdf
from .models import SlideData
from .utils import ensure_dir, first_nonempty_line


def _extract_shape_text(shape) -> list[str]:
    texts: list[str] = []
    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text)
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    texts.append(" | ".join(cells))
        except Exception:
            pass
    if getattr(shape, "shape_type", None) == 6:  # GROUP
        try:
            for sub in shape.shapes:
                texts.extend(_extract_shape_text(sub))
        except Exception:
            pass
    return texts


def extract_ppt_text(ppt_path: Path) -> list[SlideData]:
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx") from e
    prs = Presentation(str(ppt_path))
    slides: list[SlideData] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_extract_shape_text(shape))
        text = "\n".join(p.strip() for p in parts if p.strip())
        slides.append(SlideData(slide_no=idx, heading=first_nonempty_line(text), text=text, cleaned_text=text, source_type="pptx"))
    return slides


def _find_soffice() -> str | None:
    exe = shutil.which("soffice") or shutil.which("libreoffice")
    if exe:
        return exe
    candidates = [r"C:\Program Files\LibreOffice\program\soffice.exe", r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"]
    for c in candidates:
        if Path(c).exists():
            return c
    return None


def _render_ppt_with_libreoffice(ppt_path: Path, run_dir: Path) -> dict[int, Path]:
    soffice = _find_soffice()
    if not soffice:
        return {}
    tmp = ensure_dir(run_dir / "ppt_pdf")
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(tmp), str(ppt_path)], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=180)
    pdf_candidates = list(tmp.glob("*.pdf"))
    if not pdf_candidates:
        return {}
    rendered_slides = extract_pdf(pdf_candidates[0], run_dir / "ppt_render_pdf", render_scale=2.0)
    return {s.slide_no: s.image_path for s in rendered_slides if s.image_path}


def _render_ppt_with_powerpoint(ppt_path: Path, run_dir: Path) -> dict[int, Path]:
    if sys.platform != "win32":
        return {}
    try:
        import pythoncom
        import win32com.client
    except Exception:
        return {}
    out_dir = ensure_dir(run_dir / "rendered")
    try:
        pythoncom.CoInitialize()
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = 1
        pres = app.Presentations.Open(str(ppt_path.resolve()), WithWindow=False)
        pres.Export(str(out_dir.resolve()), "PNG", 1920, 1080)
        pres.Close()
        app.Quit()
        mapping: dict[int, Path] = {}
        for p in out_dir.glob("Slide*.PNG"):
            num = "".join(ch for ch in p.stem if ch.isdigit())
            if num:
                new = out_dir / f"slide_{int(num):03d}.png"
                if new.exists():
                    new.unlink()
                p.rename(new)
                mapping[int(num)] = new
        return mapping
    except Exception:
        return {}


def extract_pptx(ppt_path: Path, run_dir: Path) -> tuple[list[SlideData], list[str]]:
    warnings: list[str] = []
    slides = extract_ppt_text(ppt_path)
    image_map = _render_ppt_with_powerpoint(ppt_path, run_dir)
    if not image_map:
        try:
            image_map = _render_ppt_with_libreoffice(ppt_path, run_dir)
        except Exception as e:
            warnings.append(f"LibreOffice PPT rendering failed: {e}")
    if not image_map:
        warnings.append("PPT slide images were not rendered. Install LibreOffice or Microsoft PowerPoint for image/diagram insertion and vision extraction.")
    for s in slides:
        if s.slide_no in image_map:
            s.image_path = image_map[s.slide_no]
    return slides, warnings
